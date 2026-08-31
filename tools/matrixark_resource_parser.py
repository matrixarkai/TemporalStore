#!/usr/bin/env python3
# SPDX-License-Identifier: Apache-2.0
# Copyright 2026 MatrixArkAI
"""MatrixArk resource parsing helpers.

The parser keeps raw bytes outside TemporalStore and converts resources into
small, cited, embedding-ready chunks. The shape intentionally mirrors the useful
parts of context-filesystem systems: stable refs, section hierarchy, L0-friendly
metadata, and bounded serving payloads.
"""

from __future__ import annotations

import csv
import hashlib
import html
import json
import os
import re
import shlex
import subprocess
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any


Json = dict[str, Any]

SUPPORTED_TEXT_TYPES = {"md", "txt", "log", "html", "csv", "tsv", "json", "jsonl", "skill"}
SUPPORTED_BINARY_TYPES = {"pdf", "docx", "pptx", "xlsx", "png", "jpg", "jpeg", "webp", "gif"}
SUPPORTED_FILE_TYPES = SUPPORTED_TEXT_TYPES | SUPPORTED_BINARY_TYPES
SUPPORTED_DIRECTORY_TYPES = {"md", "txt", "pdf", "html", "csv", "tsv", "json", "jsonl", "docx", "pptx", "xlsx", "skill"}
SKIP_DIRECTORY_NAMES = {".git", ".hg", ".svn", "node_modules", "__pycache__", ".venv", "venv", "target", "build", "dist"}
DEFAULT_MAX_FILE_BYTES = int(os.environ.get("MATRIXARK_RESOURCE_MAX_FILE_BYTES", str(20 * 1024 * 1024)))
DEFAULT_MAX_DIRECTORY_FILES = int(os.environ.get("MATRIXARK_RESOURCE_MAX_DIRECTORY_FILES", "256"))
DEFAULT_MAX_DIRECTORY_DEPTH = int(os.environ.get("MATRIXARK_RESOURCE_MAX_DIRECTORY_DEPTH", "8"))
DEFAULT_MAX_TOTAL_CHUNKS = int(os.environ.get("MATRIXARK_RESOURCE_MAX_TOTAL_CHUNKS", "2048"))
DEFAULT_MAX_INLINE_TEXT_CHARS = int(os.environ.get("MATRIXARK_RESOURCE_MAX_INLINE_TEXT_CHARS", str(5 * 1024 * 1024)))
DEFAULT_TABLE_ROWS_PER_CHUNK = int(os.environ.get("MATRIXARK_RESOURCE_TABLE_ROWS_PER_CHUNK", "20"))
DEFAULT_JSON_RECORDS_PER_CHUNK = int(os.environ.get("MATRIXARK_RESOURCE_JSON_RECORDS_PER_CHUNK", "20"))
DEFAULT_MD_PACK_SECTIONS = os.environ.get("MATRIXARK_RESOURCE_MD_PACK_SECTIONS", "1") not in {"0", "false", "False", ""}
DEFAULT_SLIM_CHUNK_METADATA = os.environ.get("MATRIXARK_RESOURCE_SLIM_CHUNK_METADATA", "0") not in {"0", "false", "False", ""}
# How many tokens each encoder actually reads. A window is a property of the MODEL, not a
# constant: e5 and MiniLM stop at 512, BGE-M3 and jina-v3 at 8192. Anything beyond a model's
# limit is silently truncated by its tokenizer, and anything short of it wastes capacity and
# multiplies the number of vectors stored.
_ENCODER_WINDOWS = (
    ("bge-m3", 8192),
    ("jina-embeddings-v3", 8192),
    ("gte-multilingual", 8192),
    ("e5", 512),
    ("minilm", 512),
    ("mpnet", 384),
)
_DEFAULT_ENCODER_WINDOW = 512


def encoder_window_tokens(model: str | None = None) -> int:
    """The active encoder's token window, matched on its name."""
    name = (model or os.environ.get("MATRIXARK_EMBEDDING_MODEL")
            or os.environ.get("MATRIXARK_EMBEDDING_MODEL_PATH") or "").lower()
    for marker, window in _ENCODER_WINDOWS:
        if marker in name:
            return window
    return _DEFAULT_ENCODER_WINDOW


# The text handed to the encoder fills its window. This was 128 while chunks were 240 tokens,
# so half of every chunk never reached a vector and could only be found lexically -- through an
# index whose terms the retrieve path cannot consult.
DEFAULT_EMBEDDING_TEXT_MAX_TOKENS = int(
    os.environ.get("MATRIXARK_EMBEDDING_TEXT_MAX_TOKENS", str(encoder_window_tokens()))
)
# Chunk size is the dominant lever on ingest cost: for a 1.41 MB markdown file,
# 240-token chunks are 2126 records and 27.8 MB resident, 2000-token chunks are
# 204 records and 8.9 MB. It had no knob, so no deployment could reach it.
# Chunk size follows the encoder's window, so a chunk is never larger than what will be
# embedded. Measured on a 1 MB skill: 128-token chunks are 3,302 vectors, 256 are 1,519 and
# 512 are 744 -- and vectors are 61% of what an ingest stores, so this is the dominant lever.
DEFAULT_MAX_CHUNK_TOKENS = int(
    os.environ.get("MATRIXARK_RESOURCE_MAX_CHUNK_TOKENS", str(encoder_window_tokens()))
)
DEFAULT_OVERLAP_TOKENS = int(os.environ.get("MATRIXARK_RESOURCE_OVERLAP_TOKENS", "24"))
# Both caps bind, whichever is hit first. Deriving the character cap from the
# token cap keeps raising one from being silently cancelled by the other.
DEFAULT_MAX_CHUNK_CHARS = int(
    os.environ.get(
        "MATRIXARK_RESOURCE_MAX_CHUNK_CHARS",
        str(1400 if DEFAULT_MAX_CHUNK_TOKENS <= 240 else DEFAULT_MAX_CHUNK_TOKENS * 8),
    )
)
DEFAULT_OVERLAP_CHARS = int(os.environ.get("MATRIXARK_RESOURCE_OVERLAP_CHARS", "120"))
EMBEDDING_TEXT_PREFIX_SHARE = float(os.environ.get("MATRIXARK_EMBEDDING_TEXT_PREFIX_SHARE", "0.2"))
DEFAULT_OCR_TIMEOUT_S = float(os.environ.get("MATRIXARK_RESOURCE_OCR_TIMEOUT_S", "30"))


class ResourceParserError(RuntimeError):
    pass


def resolve_max_inline_text_chars(override: int | None = None) -> int:
    """Resolve the effective inline/text cap (chars) for resources — the single source
    of truth skills also reuse (see ``matrixark_skill_parser.resolve_max_skill_bytes``).

    Precedence: explicit ``override`` arg > ``MATRIXARK_RESOURCE_MAX_INLINE_TEXT_CHARS``
    env > ``DEFAULT_MAX_INLINE_TEXT_CHARS`` (5 MiB). Read at call time so a host-set env
    override takes effect without re-importing. Raises ``ResourceParserError`` on a
    non-positive override."""
    if override is not None:
        value = int(override)
        if value <= 0:
            raise ResourceParserError("max_inline_text_chars must be positive")
        return value
    raw = os.environ.get("MATRIXARK_RESOURCE_MAX_INLINE_TEXT_CHARS")
    if raw:
        try:
            value = int(raw)
        except ValueError:
            value = 0
        if value > 0:
            return value
    return DEFAULT_MAX_INLINE_TEXT_CHARS


@dataclass(frozen=True)
class ParsedResourceChunk:
    chunk_hash: int
    source_ref: str
    text: str
    token_estimate: int
    metadata: Json

    def to_json(self) -> Json:
        return {
            "chunk_hash": self.chunk_hash,
            "source_ref": self.source_ref,
            "text": self.text,
            "token_estimate": self.token_estimate,
            "metadata": self.metadata,
        }


def embedding_text_for_chunk(chunk: ParsedResourceChunk) -> str:
    """Return the text that should be embedded for chunk retrieval.

    The stored chunk text stays citation-friendly, while embeddings get a little
    more structure: source path, heading path, table columns, keywords, and the
    content itself. This improves query matching without bloating prompt output.
    """
    value = chunk.metadata.get("embedding_text")
    if isinstance(value, str) and value.strip():
        return value
    return build_embedding_text(chunk.text, chunk.metadata, chunk.source_ref)


def stable_hash(value: str) -> int:
    digest = hashlib.sha256(value.encode("utf-8")).digest()
    return int.from_bytes(digest[:8], "big") & 0x7FFF_FFFF_FFFF_FFFF


def infer_resource_type(raw_uri: str, resource_type: str | None = None) -> str:
    if resource_type:
        kind = resource_type.lower().lstrip(".")
    else:
        suffix = Path(raw_uri).suffix.lower().lstrip(".")
        kind = suffix or "txt"
    aliases = {
        "markdown": "md",
        "text": "txt",
        "htm": "html",
        "jpeg": "jpg",
        "jsonlines": "jsonl",
        "skill": "skill",
    }
    return aliases.get(kind, kind)


_CJK_CLASS = "぀-ヿ㐀-䶿一-鿿豈-﫿ｦ-ﾟ"
_LATIN_RUN_RE = re.compile(r"[A-Za-z0-9_]+")
_CJK_CHAR_RE = re.compile("[" + _CJK_CLASS + "]")
_OTHER_SYMBOL_RE = re.compile("[^" + _CJK_CLASS + r"\sA-Za-z0-9_]")
CJK_AWARE_TOKENS = os.environ.get("MATRIXARK_RESOURCE_CJK_TOKENS", "1") not in {"0", "false", "False", ""}
_SPLIT_SPAN_RE = re.compile(
    "[" + _CJK_CLASS + r"]|[A-Za-z0-9_]+|[^" + _CJK_CLASS + r"\sA-Za-z0-9_]+"
)


def _is_cjk_codepoint(code: int) -> bool:
    return (
        0x3040 <= code <= 0x30FF
        or 0x3400 <= code <= 0x4DBF
        or 0x4E00 <= code <= 0x9FFF
        or 0xF900 <= code <= 0xFAFF
        or 0xFF66 <= code <= 0xFF9F
    )


def _span_weight(span_text: str) -> float:
    # Spans come from _SPLIT_SPAN_RE, so the first character already determines the class:
    # a CJK span is exactly one CJK character, a latin span is all word characters, and
    # anything else is a symbol run. A codepoint test is much cheaper than re-matching the
    # span, and this runs once per span of every chunk. CJK is tested first because
    # str.isalnum() is True for CJK characters.
    code = ord(span_text[0])
    if _is_cjk_codepoint(code):
        return 0.75
    if span_text[0].isalnum() or code == 0x5F:
        return 1.1
    return 1.3 * len(span_text)


def token_estimate(text: str) -> int:
    """Estimate LLM tokens for mixed CJK/Latin text.

    Counting only ``[A-Za-z0-9_]+`` runs undercounts Chinese by ~37x (a pure-CJK
    sentence scores 1), which silently blows any ``max_context_tokens`` budget on
    a CJK corpus. Coefficients calibrated against the multilingual MiniLM
    tokenizer: median estimate/actual 1.00 over a mixed CN/EN sample.
    """
    if not CJK_AWARE_TOKENS:
        return max(1, len(_LATIN_RUN_RE.findall(text)))
    latin = len(_LATIN_RUN_RE.findall(text))
    cjk = len(_CJK_CHAR_RE.findall(text))
    other = len(_OTHER_SYMBOL_RE.findall(text))
    return max(1, int(1.1 * latin + 0.75 * cjk + 1.3 * other))


def slugify(value: str) -> str:
    value = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower())
    return value.strip("-") or "section"


def compact_ws(text: str) -> str:
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def content_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()[:16]


def normalize_parse_warnings(metadata: Json) -> list[str]:
    warnings: list[str] = []

    def add_warning(value: Any) -> None:
        text = compact_ws(str(value or ""))
        if text and text not in warnings:
            warnings.append(text[:240])

    raw_warnings = metadata.get("parse_warnings", [])
    if isinstance(raw_warnings, str):
        add_warning(raw_warnings)
    elif isinstance(raw_warnings, list):
        for value in raw_warnings:
            add_warning(value)

    if metadata.get("needs_vlm"):
        add_warning("visual_content_requires_vlm_or_ocr")
    if metadata.get("needs_ocr"):
        add_warning("scanned_or_binary_content_requires_ocr")
    if metadata.get("ocr_status") == "unavailable":
        add_warning(metadata.get("ocr_reason") or "ocr_unavailable")
    skipped = metadata.get("directory_skipped_files")
    if isinstance(skipped, int) and skipped > 0:
        add_warning(f"directory_skipped_files:{skipped}")
    return warnings


def resource_content_version(raw_uri: str, units: list[Json]) -> str:
    digest = hashlib.sha256()
    digest.update(raw_uri.encode("utf-8", errors="ignore"))
    for unit in units:
        digest.update(str(unit.get("unit_kind", "")).encode("utf-8", errors="ignore"))
        digest.update(str(unit.get("source_ref", "")).encode("utf-8", errors="ignore"))
        digest.update(str(unit.get("text", "")).encode("utf-8", errors="ignore"))
    return digest.hexdigest()[:16]


CJK_KEYWORD_BIGRAMS = os.environ.get("MATRIXARK_RESOURCE_CJK_KEYWORDS", "1") not in {"0", "false", "False", ""}
_CJK_RUN_RE = re.compile("[" + _CJK_CLASS + "]{2,}")


def keywords_for_text(text: str, limit: int = 12) -> list[str]:
    """Keywords for the secondary index.

    Latin runs alone leave Chinese text with NO keywords, so on a CJK corpus the keyword
    index — the part of the secondary index carrying any selectivity — indexes nothing.
    Chinese has no spaces to split on, so CJK runs contribute overlapping character
    bigrams, which a lexical index can match without a segmenter.

    The two are interleaved rather than concatenated. Taking Latin first and stopping at
    the limit lets a few English words in a mixed passage exhaust the quota, leaving the
    Chinese half of the same chunk unindexed — which is what happened before.
    """
    stop = {
        "the", "and", "for", "that", "with", "from", "this", "will", "are", "was", "were", "has", "have",
        "should", "into", "when", "where", "what", "which", "your", "their", "about", "after", "before",
    }
    latin: list[str] = []
    seen: set[str] = set()
    for token in re.findall(r"[A-Za-z][A-Za-z0-9_]{2,}", text.lower()):
        if token in stop or token in seen:
            continue
        seen.add(token)
        latin.append(token)
    cjk: list[str] = []
    if CJK_KEYWORD_BIGRAMS:
        for run in _CJK_RUN_RE.findall(text):
            for index in range(len(run) - 1):
                bigram = run[index : index + 2]
                if bigram in seen:
                    continue
                seen.add(bigram)
                cjk.append(bigram)
    out: list[str] = []
    for index in range(max(len(latin), len(cjk))):
        if index < len(latin):
            out.append(latin[index])
            if len(out) >= limit:
                return out
        if index < len(cjk):
            out.append(cjk[index])
            if len(out) >= limit:
                return out
    return out

def _truncate_to_tokens(text: str, max_tokens: int) -> str:
    """Trim text to roughly max_tokens estimated tokens, on a span boundary."""
    if max_tokens <= 0:
        return ""
    spans = [(m.start(), m.end()) for m in _SPLIT_SPAN_RE.finditer(text)]
    if not spans:
        return text
    used = 0.0
    for index, (start, end) in enumerate(spans):
        used += _span_weight(text[start:end])
        if used > max_tokens:
            return text[: spans[index][0]].strip()
    return text


def build_embedding_text(
    text: str,
    metadata: Json,
    source_ref: str,
    max_tokens: int | None = None,
) -> str:
    """Build the text handed to the encoder for a chunk.

    The encoder has a fixed input window (128 tokens for the multilingual MiniLM
    used here) and silently truncates past it. Emitting every metadata field
    first spent that whole window on the header: measured on markdown, the prefix
    was 123 of 128 tokens and only 2.8% of the chunk's own content reached the
    model, so chunks were matched on their heading/path/keyword header rather
    than on what they say.

    So the content leads, and a compact locating prefix is allowed only a small
    slice of the window. ``max_tokens=0`` restores the unbounded field dump.
    """
    if max_tokens is None:
        max_tokens = DEFAULT_EMBEDDING_TEXT_MAX_TOKENS
    if max_tokens <= 0:
        return _build_embedding_text_unbounded(text, metadata, source_ref)

    prefix_budget = max(8, int(max_tokens * EMBEDDING_TEXT_PREFIX_SHARE))
    locators: list[str] = []
    heading_path = metadata.get("heading_path")
    if isinstance(heading_path, list) and heading_path:
        locators.append(" / ".join(str(item) for item in heading_path if str(item).strip()))
    else:
        for key in ("heading", "title", "document_title", "relative_path"):
            value = metadata.get(key)
            if value:
                locators.append(str(value))
                break
    columns = metadata.get("columns")
    if isinstance(columns, list) and columns:
        locators.append(", ".join(str(item) for item in columns if str(item).strip()))
    prefix = _truncate_to_tokens(compact_ws(" | ".join(locators)), prefix_budget)
    body_budget = max_tokens - (token_estimate(prefix) if prefix else 0)
    body = _truncate_to_tokens(compact_ws(text), body_budget)
    return compact_ws(prefix + chr(10) + body) if prefix else compact_ws(body)


def _build_embedding_text_unbounded(text: str, metadata: Json, source_ref: str) -> str:

    fields: list[str] = []
    for key in [
        "document_title",
        "title",
        "heading",
        "section",
        "relative_path",
        "child_resource_type",
        "unit_kind",
        "resource_version",
    ]:
        value = metadata.get(key)
        if value:
            fields.append(f"{key}: {value}")
    heading_path = metadata.get("heading_path")
    if isinstance(heading_path, list) and heading_path:
        fields.append("path: " + " / ".join(str(item) for item in heading_path if str(item).strip()))
    columns = metadata.get("columns")
    if isinstance(columns, list) and columns:
        fields.append("columns: " + ", ".join(str(item) for item in columns if str(item).strip()))
    keywords = metadata.get("keywords")
    if isinstance(keywords, list) and keywords:
        fields.append("keywords: " + ", ".join(str(item) for item in keywords if str(item).strip()))
    fields.append(f"source: {source_ref}")
    fields.append(f"content: {text}")
    return compact_ws("\n".join(fields))


def summarize_resource_chunks(
    chunks: list[ParsedResourceChunk],
    *,
    raw_uri: str = "",
    resource_kind: str = "resource",
    limit: int = 900,
) -> str:
    """Build a compact L0 summary source from chunk headings and snippets."""
    parts: list[str] = []
    if raw_uri:
        parts.append(f"{resource_kind}: {raw_uri}")
    seen_paths: set[str] = set()
    for chunk in chunks[:12]:
        metadata = chunk.metadata
        heading_path = metadata.get("heading_path")
        if isinstance(heading_path, list) and heading_path:
            path_text = " / ".join(str(item) for item in heading_path if str(item).strip())
        else:
            path_text = str(metadata.get("heading") or metadata.get("section") or metadata.get("relative_path") or metadata.get("unit_kind") or "")
        if path_text and path_text not in seen_paths:
            seen_paths.add(path_text)
            parts.append(f"section: {path_text}")
        snippet = compact_ws(chunk.text)
        if snippet:
            parts.append(snippet[:360])
    return compact_ws(" ".join(parts))[:limit]


def parse_resource(
    raw_uri: str | Path,
    *,
    resource_type: str | None = None,
    text: str | None = None,
    max_chunk_chars: int = DEFAULT_MAX_CHUNK_CHARS,
    overlap_chars: int = DEFAULT_OVERLAP_CHARS,
    max_chunk_tokens: int = DEFAULT_MAX_CHUNK_TOKENS,
    overlap_tokens: int = DEFAULT_OVERLAP_TOKENS,
    chunk_hash_base: int | None = None,
    resource_version: str | None = None,
    supersedes_chunk_hashes: dict[str, int] | None = None,
    max_file_bytes: int = DEFAULT_MAX_FILE_BYTES,
    max_directory_files: int = DEFAULT_MAX_DIRECTORY_FILES,
    max_directory_depth: int = DEFAULT_MAX_DIRECTORY_DEPTH,
    max_total_chunks: int = DEFAULT_MAX_TOTAL_CHUNKS,
    max_inline_text_chars: int | None = None,
    pack_markdown_sections: bool | None = None,
    slim_chunk_metadata_fields: bool | None = None,
) -> list[ParsedResourceChunk]:
    """Parse supported resources into bounded serving chunks.

    Supported local/inline resource types: md, txt, pdf, html, csv, tsv, json,
    jsonl, docx, pptx, xlsx, image OCR/VLM, and skill. Unsupported binary files fail loudly instead of
    silently storing useless bytes.

    ``max_inline_text_chars`` defaults (None) to ``resolve_max_inline_text_chars()`` —
    the shared inline/text cap (5 MiB, knob ``MATRIXARK_RESOURCE_MAX_INLINE_TEXT_CHARS``)
    that skills reuse too, so a skill and a resource of equal size gate identically.
    """
    max_inline_text_chars = resolve_max_inline_text_chars(max_inline_text_chars)
    if max_chunk_chars < 256:
        raise ResourceParserError("max_chunk_chars must be >= 256")
    if overlap_chars < 0 or overlap_chars >= max_chunk_chars:
        raise ResourceParserError("overlap_chars must be non-negative and smaller than max_chunk_chars")
    if max_chunk_tokens < 32:
        raise ResourceParserError("max_chunk_tokens must be >= 32")
    if overlap_tokens < 0 or overlap_tokens >= max_chunk_tokens:
        raise ResourceParserError("overlap_tokens must be non-negative and smaller than max_chunk_tokens")
    if max_total_chunks <= 0:
        raise ResourceParserError("max_total_chunks must be positive")
    if max_inline_text_chars <= 0:
        raise ResourceParserError("max_inline_text_chars must be positive")
    if text is not None and len(text) > max_inline_text_chars:
        raise ResourceParserError(f"inline resource text too large: {len(text)} chars, max is {max_inline_text_chars}")
    raw_uri_text = str(raw_uri)
    kind = infer_resource_type(raw_uri_text, resource_type)
    units = _parse_units(
        raw_uri_text,
        kind,
        text,
        max_file_bytes=max_file_bytes,
        max_directory_files=max_directory_files,
        max_directory_depth=max_directory_depth,
        max_inline_text_chars=max_inline_text_chars,
    )

    if slim_chunk_metadata_fields is None:
        slim_chunk_metadata_fields = DEFAULT_SLIM_CHUNK_METADATA
    if pack_markdown_sections is None:
        pack_markdown_sections = DEFAULT_MD_PACK_SECTIONS
    if pack_markdown_sections and kind in {"md", "skill"}:
        units = _pack_markdown_units(units, max_chunk_tokens)

    chunks: list[ParsedResourceChunk] = []
    version = resource_version or resource_content_version(raw_uri_text, units)
    supersedes = supersedes_chunk_hashes or {}
    for unit_index, unit in enumerate(units):
        unit_text = compact_ws(str(unit.get("text", "")))
        for split_index, piece in enumerate(
            _split_text(
                unit_text,
                max_chunk_chars,
                overlap_chars,
                max_chunk_tokens=max_chunk_tokens,
                overlap_tokens=overlap_tokens,
            )
        ):
            piece = compact_ws(piece)
            if not piece:
                continue
            piece_hash = content_hash(piece)
            piece_tokens = token_estimate(piece)
            metadata = {key: value for key, value in unit.items() if key != "text"}
            metadata.update(
                {
                    "resource_type": kind,
                    "resource_version": version,
                    "chunk_index": len(chunks),
                    "unit_index": unit_index,
                    "split_index": split_index,
                    "char_count": len(piece),
                    "token_count": piece_tokens,
                    "splitter": "token_aware_v1",
                    "max_chunk_tokens": max_chunk_tokens,
                    "overlap_tokens": overlap_tokens,
                    "content_hash": piece_hash,
                    "keywords": [] if slim_chunk_metadata_fields else keywords_for_text(piece),
                }
            )
            if "page" in metadata and split_index:
                metadata["page_section"] = split_index
            source_ref = _source_ref(raw_uri_text, metadata)
            metadata["citation"] = source_ref
            metadata["supersedes_chunk_hash"] = supersedes.get(source_ref) or supersedes.get(piece_hash)
            # Only build the encoder string when it will actually be kept. It is 41% of
            # parse time and slim metadata drops it, so computing it there is pure waste;
            # the embedding step recomputes it from the chunk when it needs it.
            if not slim_chunk_metadata_fields:
                metadata["embedding_text"] = build_embedding_text(piece, metadata, source_ref)
            metadata["parse_warnings"] = normalize_parse_warnings(metadata)
            metadata["raw_storage_policy"] = "raw_uri_only"
            metadata["raw_bytes_stored"] = False
            if slim_chunk_metadata_fields:
                metadata = slim_chunk_metadata(metadata)
            chunk_hash = (
                chunk_hash_base + len(chunks)
                if chunk_hash_base is not None
                else stable_hash(f"resource_chunk:{source_ref}:{version}:{piece_hash}")
            )
            if len(chunks) >= max_total_chunks:
                raise ResourceParserError(f"resource produced more than max_total_chunks={max_total_chunks}")
            chunks.append(
                ParsedResourceChunk(
                    chunk_hash=chunk_hash,
                    source_ref=source_ref,
                    text=piece,
                    token_estimate=piece_tokens,
                    metadata=metadata,
                )
            )
    return chunks


def _parse_units(
    raw_uri: str,
    kind: str,
    text: str | None,
    *,
    max_file_bytes: int = DEFAULT_MAX_FILE_BYTES,
    max_directory_files: int = DEFAULT_MAX_DIRECTORY_FILES,
    max_directory_depth: int = DEFAULT_MAX_DIRECTORY_DEPTH,
    max_inline_text_chars: int = DEFAULT_MAX_INLINE_TEXT_CHARS,
) -> list[Json]:
    if text is None:
        path = Path(raw_uri)
        if not path.exists():
            raise ResourceParserError(f"resource path does not exist: {raw_uri}")
        if path.is_dir():
            return _parse_directory_units(path, max_file_bytes=max_file_bytes, max_directory_files=max_directory_files, max_directory_depth=max_directory_depth)
        if kind not in SUPPORTED_FILE_TYPES:
            raise ResourceParserError(f"unsupported resource type for local file: {kind}")
        _ensure_file_size(path, max_file_bytes)
        if kind == "pdf":
            return _parse_pdf_units(path, raw_uri)
        if kind == "docx":
            return _parse_docx_units(path, raw_uri)
        if kind == "pptx":
            return _parse_pptx_units(path, raw_uri)
        if kind == "xlsx":
            return _parse_xlsx_units(path, raw_uri)
        if kind in {"png", "jpg", "jpeg", "webp", "gif"}:
            return _image_units(path, raw_uri, kind)
        text = _read_text_file(path, max_file_bytes)
    if text is not None and len(text) > max_inline_text_chars:
        raise ResourceParserError(f"inline resource text too large: {len(text)} chars, max is {max_inline_text_chars}")
    if kind in {"md", "skill"}:
        return _markdown_units(text)
    if kind == "html":
        return _html_units(text, raw_uri)
    if kind in {"csv", "tsv"}:
        return _delimited_units(text, raw_uri, delimiter="\t" if kind == "tsv" else ",")
    if kind == "jsonl":
        return _jsonl_units(text, raw_uri)
    if kind == "json":
        return _json_units(text, raw_uri)
    if kind in {"txt", "log"}:
        return _paragraph_units(text, raw_uri)
    if kind in {"png", "jpg", "jpeg", "webp", "gif"}:
        return _binary_stub_units(raw_uri, kind, reason="inline image bytes are not accepted; pass a local image path for OCR/VLM parsing")
    return _paragraph_units(text, raw_uri)


def _ensure_file_size(path: Path, max_file_bytes: int) -> None:
    try:
        size = path.stat().st_size
    except OSError as exc:
        raise ResourceParserError(f"cannot stat resource file {path}: {exc}") from exc
    if size > max_file_bytes:
        raise ResourceParserError(f"resource file too large: {path} has {size} bytes, max is {max_file_bytes}")


def _read_text_file(path: Path, max_file_bytes: int) -> str:
    _ensure_file_size(path, max_file_bytes)
    try:
        return path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8-sig", errors="replace")


def _markdown_units(text: str) -> list[Json]:
    document_metadata, text = _split_simple_front_matter(text)
    document_title = ""
    units: list[Json] = []
    heading_stack: list[tuple[int, str]] = []
    current_heading = "document"
    current_level = 0
    current_start_line = 1
    buffer: list[str] = []
    in_code = False

    def path_slugs() -> list[str]:
        return [slugify(name) for _, name in heading_stack] or ["document"]

    def flush(end_line: int) -> None:
        content = "\n".join(buffer).strip()
        if content:
            units.append(
                {
                    "text": content,
                    "heading": current_heading,
                    "heading_slug": slugify(current_heading),
                    "heading_level": current_level,
                    "heading_path": [name for _, name in heading_stack] or ["document"],
                    "heading_path_slugs": path_slugs(),
                    "line_start": current_start_line,
                    "line_end": max(current_start_line, end_line),
                    "unit_kind": "markdown_section",
                    "document_metadata": document_metadata,
                    "document_title": document_title or (heading_stack[0][1] if heading_stack else ""),
                }
            )
        buffer.clear()

    for line_no, line in enumerate(text.splitlines(), start=1):
        if line.strip().startswith("```") or line.strip().startswith("~~~"):
            in_code = not in_code
            buffer.append(line.rstrip())
            continue
        match = None if in_code else re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if match:
            flush(line_no - 1)
            current_level = len(match.group(1))
            current_heading = match.group(2).strip()
            heading_stack = [item for item in heading_stack if item[0] < current_level]
            heading_stack.append((current_level, current_heading))
            if current_level == 1 and not document_title:
                document_title = current_heading
            current_start_line = line_no
            buffer.append(line.strip())
        else:
            buffer.append(line.rstrip())
    flush(len(text.splitlines()) or 1)
    if not units and text.strip():
        units.append(
            {
                "text": text.strip(),
                "heading": "document",
                "heading_slug": "document",
                "heading_level": 0,
                "heading_path": ["document"],
                "heading_path_slugs": ["document"],
                "line_start": 1,
                "line_end": len(text.splitlines()) or 1,
                "unit_kind": "markdown_document",
                "document_metadata": document_metadata,
                "document_title": document_title,
            }
        )
    return units


def _pack_markdown_units(units: list[Json], max_chunk_tokens: int) -> list[Json]:
    """Merge consecutive small markdown sections up to the chunk token budget.

    ``_markdown_units`` emits one unit per heading, so a document made of many
    short sections produces many chunks far below ``max_chunk_tokens``. Packing
    adjacent siblings keeps each chunk near the budget without splitting a
    section across chunks. The pack keeps the first section's heading and
    records every heading it covers, so citations stay resolvable.
    """
    if not units:
        return units
    packed: list[Json] = []
    current: Json | None = None
    current_tokens = 0
    joiner = chr(10) + chr(10)
    for unit in units:
        if unit.get("unit_kind") != "markdown_section":
            if current is not None:
                packed.append(current)
                current = None
                current_tokens = 0
            packed.append(unit)
            continue
        tokens = token_estimate(str(unit.get("text", "")))
        if tokens >= max_chunk_tokens:
            if current is not None:
                packed.append(current)
                current = None
                current_tokens = 0
            packed.append(unit)
            continue
        parent = tuple(unit.get("heading_path", [])[:-1])
        if current is not None:
            same_parent = tuple(current.get("_pack_parent", ())) == parent
            if same_parent and current_tokens + tokens <= max_chunk_tokens:
                current["text"] = current["text"] + joiner + str(unit.get("text", ""))
                current["line_end"] = unit.get("line_end", current.get("line_end"))
                current["packed_headings"].append(unit.get("heading", ""))
                current["packed_sections"] = len(current["packed_headings"])
                current_tokens += tokens
                continue
            packed.append(current)
        current = dict(unit)
        current["_pack_parent"] = parent
        current["packed_headings"] = [unit.get("heading", "")]
        current["packed_sections"] = 1
        current_tokens = tokens
    if current is not None:
        packed.append(current)
    for unit in packed:
        unit.pop("_pack_parent", None)
    return packed




# Fields a served chunk actually needs: enough to retrieve it, cite it, order it
# within its document, and detect that its content changed. Everything else is
# recomputable from the text or from the resource manifest.
SLIM_CHUNK_METADATA_FIELDS = frozenset({
    "resource_type",
    "resource_version",
    "chunk_index",
    "unit_index",
    "split_index",
    "content_hash",
    "citation",
    "token_count",
    "heading",
    "heading_slug",
    "line_start",
    "line_end",
    "page",
    "page_number",
    "slide_number",
    "record_index",
    "record_start",
    "record_end",
    "row_index",
    "row_range",
    "unit_kind",
    "packed_sections",
    "parse_warnings",
    "supersedes_chunk_hash",
})


def slim_chunk_metadata(metadata: Json) -> Json:
    """Drop metadata that is duplicated, derivable, or constant.

    ``embedding_text`` alone is ~67% of a chunk's metadata: it is a second copy
    of the chunk text enriched for the encoder, and it is only needed while the
    vector is being produced. ``keywords`` feeds a lexical posting list that a
    vector index already covers. Both are recomputed on demand via
    ``build_embedding_text`` / ``keywords_for_text`` when a caller wants them.
    """
    return {key: value for key, value in metadata.items() if key in SLIM_CHUNK_METADATA_FIELDS}


def _split_simple_front_matter(text: str) -> tuple[Json, str]:
    if not text.startswith("---"):
        return {}, text
    match = re.match(r"^---\s*\n(.*?)\n---\s*\n?(.*)$", text, flags=re.DOTALL)
    if not match:
        return {}, text
    metadata: Json = {}
    current_key = ""
    for raw_line in match.group(1).splitlines():
        stripped = raw_line.strip()
        if not stripped or stripped.startswith("#"):
            continue
        if stripped.startswith("-") and current_key:
            metadata.setdefault(current_key, [])
            if isinstance(metadata[current_key], list):
                metadata[current_key].append(stripped[1:].strip().strip('"').strip("'"))
            continue
        if ":" not in raw_line:
            continue
        key, value = raw_line.split(":", 1)
        current_key = key.strip()
        value = value.strip()
        if not value:
            metadata[current_key] = []
        elif value.startswith("[") and value.endswith("]"):
            metadata[current_key] = [item.strip().strip('"').strip("'") for item in value[1:-1].split(",") if item.strip()]
        else:
            metadata[current_key] = value.strip('"').strip("'")
    return metadata, match.group(2)


def _paragraph_units(text: str, raw_uri: str) -> list[Json]:
    units: list[Json] = []
    cursor_line = 1
    for index, part in enumerate(re.split(r"\n\s*\n+", text)):
        paragraph = part.strip()
        line_count = max(1, part.count("\n") + 1)
        if paragraph:
            units.append(
                {
                    "text": paragraph,
                    "paragraph_index": index,
                    "line_start": cursor_line,
                    "line_end": cursor_line + line_count - 1,
                    "section": Path(raw_uri).name or "document",
                    "unit_kind": "paragraph",
                }
            )
        cursor_line += line_count
    if not units and text.strip():
        units.append({"text": text.strip(), "paragraph_index": 0, "section": Path(raw_uri).name or "document", "unit_kind": "paragraph"})
    return units


class _TextHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.title = ""
        self._in_title = False

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"p", "div", "section", "article", "br", "li", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")
        if tag == "title":
            self._in_title = True

    def handle_endtag(self, tag: str) -> None:
        if tag in {"p", "div", "section", "article", "li", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")
        if tag == "title":
            self._in_title = False

    def handle_data(self, data: str) -> None:
        value = html.unescape(data).strip()
        if not value:
            return
        if self._in_title:
            self.title = value
        self.parts.append(value + " ")


def _html_units(text: str, raw_uri: str) -> list[Json]:
    parser = _TextHTMLParser()
    parser.feed(text)
    plain = compact_ws("".join(parser.parts))
    units = _paragraph_units(plain, raw_uri)
    for unit in units:
        unit["unit_kind"] = "html_text"
        if parser.title:
            unit["title"] = parser.title
    return units


def _delimited_units(text: str, raw_uri: str, *, delimiter: str) -> list[Json]:
    rows = list(csv.DictReader(text.splitlines(), delimiter=delimiter))
    if not rows:
        return _paragraph_units(text, raw_uri)
    units: list[Json] = []
    chunk_size = max(1, DEFAULT_TABLE_ROWS_PER_CHUNK)
    for start in range(0, len(rows), chunk_size):
        group = rows[start : start + chunk_size]
        rendered_rows: list[str] = []
        columns: list[str] = []
        for row_offset, row in enumerate(group):
            clean = {str(k): str(v) for k, v in row.items() if k is not None and str(v).strip()}
            if not columns:
                columns = list(clean.keys())
            if clean:
                rendered_rows.append(f"row {start + row_offset}: " + "; ".join(f"{key}: {value}" for key, value in clean.items()))
        if rendered_rows:
            row_end = start + len(group) - 1
            unit: Json = {
                "text": "\n".join(rendered_rows),
                "row_start": start,
                "row_end": row_end,
                "row_count": len(group),
                "columns": columns,
                "unit_kind": "table_row_group" if len(group) > 1 else "table_row",
            }
            if len(group) == 1:
                unit["row_index"] = start
            units.append(unit)
    return units or _paragraph_units(text, raw_uri)


def _jsonl_units(text: str, raw_uri: str) -> list[Json]:
    parsed: list[tuple[int, str, str]] = []
    for index, line in enumerate(text.splitlines()):
        if not line.strip():
            continue
        try:
            obj = json.loads(line)
            rendered = json.dumps(obj, ensure_ascii=False, sort_keys=True)
            parsed.append((index, rendered, "jsonl_record"))
        except json.JSONDecodeError:
            parsed.append((index, line.strip(), "jsonl_text"))
    return _record_group_units(parsed, "jsonl") or _paragraph_units(text, raw_uri)


def _json_units(text: str, raw_uri: str) -> list[Json]:
    try:
        obj = json.loads(text)
    except json.JSONDecodeError:
        return _paragraph_units(text, raw_uri)
    if isinstance(obj, list):
        parsed = [(index, json.dumps(item, ensure_ascii=False, sort_keys=True), "json_record") for index, item in enumerate(obj)]
        return _record_group_units(parsed, "json") or _paragraph_units(text, raw_uri)
    if isinstance(obj, dict):
        return [{"text": json.dumps(obj, ensure_ascii=False, sort_keys=True), "record_index": 0, "unit_kind": "json_document"}]
    return [{"text": str(obj), "record_index": 0, "unit_kind": "json_value"}]


def _record_group_units(records: list[tuple[int, str, str]], family: str) -> list[Json]:
    units: list[Json] = []
    chunk_size = max(1, DEFAULT_JSON_RECORDS_PER_CHUNK)
    for start in range(0, len(records), chunk_size):
        group = records[start : start + chunk_size]
        if not group:
            continue
        record_start = group[0][0]
        record_end = group[-1][0]
        rendered = "\n".join(f"record {index}: {text}" for index, text, _kind in group if text)
        if not rendered.strip():
            continue
        unit: Json = {
            "text": rendered,
            "record_start": record_start,
            "record_end": record_end,
            "record_count": len(group),
            "unit_kind": f"{family}_record_group" if len(group) > 1 else group[0][2],
        }
        if len(group) == 1:
            unit["record_index"] = record_start
        units.append(unit)
    return units


def _binary_stub_units(raw_uri: str, kind: str, *, reason: str = "") -> list[Json]:
    warning = reason or "no OCR/VLM provider returned text"
    return [
        {
            "text": f"Binary {kind} resource at {raw_uri}. Use a VLM or OCR provider to extract visual content.",
            "unit_kind": "binary_stub",
            "needs_vlm": True,
            "ocr_status": "unavailable",
            "ocr_reason": warning,
            "parse_warnings": [warning, "raw_bytes_not_stored"],
        }
    ]


def _image_units(path: Path, raw_uri: str, kind: str) -> list[Json]:
    ocr = _ocr_path_text(path, kind=kind)
    if ocr.get("text"):
        return [
            {
                "text": ocr["text"],
                "unit_kind": "image_vlm" if ocr.get("provider") == "vlm_command" else "image_ocr",
                "ocr_provider": ocr.get("provider", ""),
                "ocr_status": "ok",
                "image_index": 0,
            }
        ]
    units = _binary_stub_units(raw_uri, kind, reason=str(ocr.get("error") or "OCR/VLM provider not configured"))
    units[0]["unit_kind"] = "image_ocr_stub"
    units[0]["image_index"] = 0
    return units


def _parse_directory_units(
    path: Path,
    *,
    max_file_bytes: int,
    max_directory_files: int,
    max_directory_depth: int,
) -> list[Json]:
    units: list[Json] = []
    parsed_files = 0
    skipped_files = 0
    for child in sorted(item for item in path.rglob("*") if item.is_file()):
        relative = child.relative_to(path)
        if child.is_symlink():
            skipped_files += 1
            continue
        if len(relative.parts) > max_directory_depth:
            skipped_files += 1
            continue
        if any(part.startswith(".") or part in SKIP_DIRECTORY_NAMES for part in relative.parts):
            skipped_files += 1
            continue
        child_kind = infer_resource_type(str(child))
        if child_kind not in SUPPORTED_DIRECTORY_TYPES:
            skipped_files += 1
            continue
        if parsed_files >= max_directory_files:
            raise ResourceParserError(f"directory resource has more than max_directory_files={max_directory_files}: {path}")
        try:
            child_units = _parse_units(str(child), child_kind, None, max_file_bytes=max_file_bytes, max_directory_files=max_directory_files, max_directory_depth=max_directory_depth)
        except ResourceParserError:
            skipped_files += 1
            continue
        parsed_files += 1
        relative_path = str(relative).replace("\\", "/")
        for unit in child_units:
            unit = dict(unit)
            unit["child_uri"] = str(child)
            unit["relative_path"] = relative_path
            unit["child_resource_type"] = child_kind
            unit["directory_root"] = str(path)
            units.append(unit)
    if not units:
        raise ResourceParserError(f"directory resource has no supported files: {path}")
    for unit in units:
        unit["directory_file_count"] = parsed_files
        unit["directory_skipped_files"] = skipped_files
        if skipped_files > 0:
            warnings = unit.setdefault("parse_warnings", [])
            if isinstance(warnings, list):
                warnings.append(f"directory_skipped_files:{skipped_files}")
    return units


def _parse_pdf_units(path: Path, raw_uri: str) -> list[Json]:
    data = path.read_bytes()
    if not data.startswith(b"%PDF"):
        units = _paragraph_units(data.decode("utf-8", errors="replace"), raw_uri)
        for unit in units:
            unit["unit_kind"] = "pdf_text_fallback"
        return units
    errors: list[str] = []
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        units = []
        for index, page in enumerate(reader.pages):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                units.append({"text": page_text, "page": index + 1, "unit_kind": "pdf_page"})
        if units:
            return units
    except Exception as exc:  # pragma: no cover - optional dependency path.
        errors.append(f"pypdf: {exc}")
    try:
        import pdfplumber  # type: ignore

        units = []
        with pdfplumber.open(str(path)) as pdf:
            for index, page in enumerate(pdf.pages):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    units.append({"text": page_text, "page": index + 1, "unit_kind": "pdf_page"})
        if units:
            return units
    except Exception as exc:  # pragma: no cover - optional dependency path.
        errors.append(f"pdfplumber: {exc}")
    ocr_units = _ocr_pdf_units(path, raw_uri)
    if ocr_units:
        return ocr_units
    stub = _binary_stub_units(
        raw_uri,
        "pdf",
        reason="; ".join(errors) or "PDF has no extractable text and OCR/VLM provider is not configured",
    )
    stub[0]["unit_kind"] = "pdf_ocr_stub"
    stub[0]["needs_ocr"] = True
    return stub


def _ocr_pdf_units(path: Path, raw_uri: str) -> list[Json]:
    whole_pdf = _ocr_path_text(path, kind="pdf")
    if whole_pdf.get("text"):
        units = _paragraph_units(str(whole_pdf["text"]), raw_uri)
        for unit in units:
            unit["unit_kind"] = "pdf_ocr"
            unit["ocr_provider"] = whole_pdf.get("provider", "")
            unit["ocr_status"] = "ok"
        return units
    try:
        from pdf2image import convert_from_path  # type: ignore
        import pytesseract  # type: ignore
    except Exception:
        return []
    units: list[Json] = []
    try:
        for index, image in enumerate(convert_from_path(str(path)), start=1):
            page_text = compact_ws(pytesseract.image_to_string(image) or "")
            if page_text:
                units.append(
                    {
                        "text": page_text,
                        "page": index,
                        "unit_kind": "pdf_page_ocr",
                        "ocr_provider": "pytesseract+pdf2image",
                        "ocr_status": "ok",
                    }
                )
    except Exception:
        return []
    return units


def _ocr_path_text(path: Path, *, kind: str) -> Json:
    vlm_command = os.environ.get("MATRIXARK_RESOURCE_VLM_COMMAND", "").strip()
    if vlm_command:
        argv = [part.format(path=str(path), kind=kind) for part in shlex.split(vlm_command)]
        if not argv:
            return {"text": "", "provider": "vlm_command", "error": "empty VLM command"}
        try:
            result = subprocess.run(
                argv,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=DEFAULT_OCR_TIMEOUT_S,
                check=False,
            )
        except Exception as exc:
            return {"text": "", "provider": "vlm_command", "error": str(exc)}
        if result.returncode == 0 and result.stdout.strip():
            return {"text": compact_ws(result.stdout), "provider": "vlm_command"}
        return {
            "text": "",
            "provider": "vlm_command",
            "error": compact_ws(result.stderr or result.stdout or f"VLM command exited {result.returncode}"),
        }
    command = os.environ.get("MATRIXARK_RESOURCE_OCR_COMMAND", "").strip()
    if command:
        argv = [part.format(path=str(path), kind=kind) for part in shlex.split(command)]
        if not argv:
            return {"text": "", "provider": "command", "error": "empty OCR command"}
        try:
            result = subprocess.run(
                argv,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=DEFAULT_OCR_TIMEOUT_S,
                check=False,
            )
        except Exception as exc:
            return {"text": "", "provider": "command", "error": str(exc)}
        if result.returncode == 0 and result.stdout.strip():
            return {"text": compact_ws(result.stdout), "provider": "command"}
        return {
            "text": "",
            "provider": "command",
            "error": compact_ws(result.stderr or result.stdout or f"OCR command exited {result.returncode}"),
        }
    try:
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore
    except Exception as exc:
        return {"text": "", "provider": "pytesseract", "error": f"pytesseract/Pillow unavailable: {exc}"}
    try:
        text = compact_ws(pytesseract.image_to_string(Image.open(path)) or "")
    except Exception as exc:
        return {"text": "", "provider": "pytesseract", "error": str(exc)}
    return {"text": text, "provider": "pytesseract", "error": "" if text else "OCR produced no text"}


def _parse_pptx_units(path: Path, raw_uri: str) -> list[Json]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency path.
        raise ResourceParserError(f"Unable to parse PPTX resource. Install python-pptx: {exc}")
    presentation = Presentation(str(path))
    units: list[Json] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                parts.append(text.strip())
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                parts.append("speaker notes: " + notes.strip())
        except Exception:
            pass
        text = compact_ws("\n".join(parts))
        if text:
            title = parts[0].splitlines()[0][:120] if parts else f"slide {slide_index}"
            units.append(
                {
                    "text": text,
                    "slide_number": slide_index,
                    "heading": title,
                    "heading_slug": slugify(title),
                    "unit_kind": "pptx_slide",
                }
            )
    if not units:
        raise ResourceParserError(f"PPTX resource has no extractable slide text: {raw_uri}")
    return units


def _parse_xlsx_units(path: Path, raw_uri: str) -> list[Json]:
    try:
        import openpyxl  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency path.
        raise ResourceParserError(f"Unable to parse XLSX resource. Install openpyxl: {exc}")
    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    units: list[Json] = []
    chunk_size = max(1, DEFAULT_TABLE_ROWS_PER_CHUNK)
    for worksheet in workbook.worksheets:
        rows_iter = worksheet.iter_rows(values_only=True)
        try:
            header_row = next(rows_iter)
        except StopIteration:
            continue
        columns = [str(value).strip() if value is not None and str(value).strip() else f"column_{idx + 1}" for idx, value in enumerate(header_row)]
        row_buffer: list[tuple[int, list[Any]]] = []

        def flush() -> None:
            if not row_buffer:
                return
            rendered_rows: list[str] = []
            for row_number, row_values in row_buffer:
                clean = []
                for idx, value in enumerate(row_values):
                    if value is None or not str(value).strip():
                        continue
                    key = columns[idx] if idx < len(columns) else f"column_{idx + 1}"
                    clean.append(f"{key}: {value}")
                if clean:
                    rendered_rows.append(f"row {row_number}: " + "; ".join(clean))
            if rendered_rows:
                units.append(
                    {
                        "text": "\n".join(rendered_rows),
                        "sheet_name": worksheet.title,
                        "row_start": row_buffer[0][0],
                        "row_end": row_buffer[-1][0],
                        "row_count": len(row_buffer),
                        "columns": columns,
                        "unit_kind": "xlsx_row_group" if len(row_buffer) > 1 else "xlsx_row",
                    }
                )
            row_buffer.clear()

        for row_number, row_values in enumerate(rows_iter, start=2):
            if not any(value is not None and str(value).strip() for value in row_values):
                continue
            row_buffer.append((row_number, list(row_values)))
            if len(row_buffer) >= chunk_size:
                flush()
        flush()
    if not units:
        raise ResourceParserError(f"XLSX resource has no extractable table rows: {raw_uri}")
    return units


def _parse_docx_units(path: Path, raw_uri: str) -> list[Json]:
    try:
        import docx  # type: ignore
    except Exception as exc:  # pragma: no cover - optional dependency path.
        raise ResourceParserError(f"Unable to parse DOCX resource. Install python-docx: {exc}")
    document = docx.Document(str(path))
    units = []
    current_heading = "document"
    buffer: list[str] = []

    def flush(index: int) -> None:
        content = "\n".join(buffer).strip()
        if content:
            units.append({"text": content, "heading": current_heading, "heading_slug": slugify(current_heading), "paragraph_index": index, "unit_kind": "docx_section"})
        buffer.clear()

    for index, paragraph in enumerate(document.paragraphs):
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        text = paragraph.text.strip()
        if not text:
            continue
        if style.startswith("heading"):
            flush(index)
            current_heading = text
            buffer.append(text)
        else:
            buffer.append(text)
    flush(len(document.paragraphs))
    return units or _paragraph_units("\n\n".join(p.text for p in document.paragraphs), raw_uri)


def _split_text(
    text: str,
    max_chunk_chars: int,
    overlap_chars: int,
    *,
    max_chunk_tokens: int,
    overlap_tokens: int,
) -> list[str]:
    text = compact_ws(text)
    if len(text) <= max_chunk_chars and token_estimate(text) <= max_chunk_tokens:
        return [text]

    token_pieces = _split_text_by_tokens(text, max_chunk_tokens, overlap_tokens)
    pieces: list[str] = []
    for token_piece in token_pieces:
        if len(token_piece) <= max_chunk_chars:
            pieces.append(token_piece)
            continue
        pieces.extend(_split_text_by_chars(token_piece, max_chunk_chars, overlap_chars))
    return pieces


def _split_text_by_tokens(text: str, max_chunk_tokens: int, overlap_tokens: int) -> list[str]:
    spans = [(match.start(), match.end()) for match in _SPLIT_SPAN_RE.finditer(text)]
    # Weight each span the way token_estimate does, so a piece that fits
    # max_chunk_tokens spans also fits max_chunk_tokens estimated tokens.
    weights = [_span_weight(text[a:b]) for a, b in spans]
    cumulative = [0.0]
    for w in weights:
        cumulative.append(cumulative[-1] + w)
    if cumulative[-1] <= max_chunk_tokens:
        return [text]
    pieces: list[str] = []
    start_token = 0
    while start_token < len(spans):
        budget = cumulative[start_token] + max_chunk_tokens
        end_token = start_token + 1
        while end_token < len(spans) and cumulative[end_token + 1] <= budget:
            end_token += 1
        end_token = min(len(spans), max(end_token, start_token + 1))
        if end_token < len(spans):
            boundary = _token_boundary(text, spans, start_token, end_token)
            if cumulative[boundary] - cumulative[start_token] > max_chunk_tokens / 2:
                end_token = boundary
        start_char = spans[start_token][0]
        end_char = spans[end_token - 1][1]
        piece = text[start_char:end_char].strip()
        if piece:
            pieces.append(piece)
        if end_token >= len(spans):
            break
        start_token = max(end_token - overlap_tokens, start_token + 1)
    return pieces


def _token_boundary(text: str, spans: list[tuple[int, int]], start_token: int, end_token: int) -> int:
    for index in range(end_token - 1, start_token, -1):
        token_end = spans[index - 1][1]
        if token_end < len(text) and text[token_end : token_end + 1] in {".", "!", "?", "\n", ";"}:
            return index
    return end_token


def _split_text_by_chars(text: str, max_chunk_chars: int, overlap_chars: int) -> list[str]:
    pieces: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + max_chunk_chars)
        if end < len(text):
            boundary = max(text.rfind("\n\n", start, end), text.rfind(". ", start, end), text.rfind("; ", start, end))
            if boundary > start + max_chunk_chars // 2:
                end = boundary + 1
        piece = text[start:end].strip()
        if piece:
            pieces.append(piece)
        if end >= len(text):
            break
        start = max(end - overlap_chars, start + 1)
    return pieces


def _source_ref(raw_uri: str, metadata: Json) -> str:
    base_uri = str(metadata.get("child_uri") or raw_uri)
    if "page" in metadata:
        suffix = f"page={metadata['page']}"
        if metadata.get("split_index", 0):
            suffix += f"&section={metadata['split_index']}"
    elif "heading_path_slugs" in metadata:
        suffix = "heading=" + "/".join(str(part) for part in metadata["heading_path_slugs"])
    elif "heading_slug" in metadata:
        suffix = f"heading={metadata['heading_slug']}"
    elif "row_start" in metadata and "row_end" in metadata:
        suffix = (
            f"row={metadata['row_start']}"
            if metadata["row_start"] == metadata["row_end"]
            else f"row={metadata['row_start']}-{metadata['row_end']}"
        )
    elif "row_index" in metadata:
        suffix = f"row={metadata['row_index']}"
    elif "record_start" in metadata and "record_end" in metadata:
        suffix = (
            f"record={metadata['record_start']}"
            if metadata["record_start"] == metadata["record_end"]
            else f"record={metadata['record_start']}-{metadata['record_end']}"
        )
    elif "record_index" in metadata:
        suffix = f"record={metadata['record_index']}"
    elif "paragraph_index" in metadata:
        suffix = f"paragraph={metadata['paragraph_index']}"
    else:
        suffix = f"chunk={metadata.get('chunk_index', 0)}"
    if metadata.get("split_index", 0) and "page" not in metadata:
        suffix += f"&part={metadata['split_index']}"
    return f"{base_uri}#{suffix}"
